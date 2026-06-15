from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Laporan Vulnerability Assessment DVWA"
    subtitle.text = "Analisis Keamanan Aplikasi Web\n\nBerdasarkan Laporan Praktikum"

    # Helper function to add bullet slide
    def add_bullet_slide(prs, title_text, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = title_text
        body = slide.placeholders[1]
        tf = body.text_frame
        for i, point in enumerate(bullet_points):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        return slide

    # Slide 2: Pendahuluan
    add_bullet_slide(prs, "Pendahuluan", [
        "Keamanan aplikasi web sangat penting karena mengelola informasi sensitif.",
        "Ancaman umum: SQL Injection, XSS, Command Injection, File Upload.",
        "Target Pengujian: DVWA (Damn Vulnerable Web Application) v1.10.",
        "Lingkungan: Terisolasi menggunakan Docker di dalam VM Kali Linux.",
        "Tujuan: Mengidentifikasi kerentanan dan menyusun rekomendasi mitigasi."
    ])

    # Slide 3: Ruang Lingkup & Metodologi
    add_bullet_slide(prs, "Ruang Lingkup & Metodologi", [
        "Target: DVWA v1.10 (http://localhost)",
        "Platform: Docker Container (vulnerables/web-dvwa)",
        "Metode Pengujian: Black-Box Testing",
        "Fase Pengujian:",
        "  1. Persiapan (Setup Lingkungan)",
        "  2. Reconnaissance (Nmap, WhatWeb, Nikto, Gobuster)",
        "  3. Scanning (OWASP ZAP)",
        "  4. Exploitation (Burp Suite, Manual)",
        "  5. Reporting (CVSS Scoring, Rekomendasi)"
    ])

    # Slide 4: Persiapan & Konfigurasi Lingkungan
    add_bullet_slide(prs, "Persiapan & Konfigurasi Lingkungan", [
        "Arsitektur Lingkungan:",
        "  - Laptop/PC Host → VirtualBox (NAT) → Kali Linux VM",
        "  - Kali Linux → Docker Engine → Container DVWA (port 80)",
        "Konfigurasi DVWA:",
        "  - Security Level: Low (Tidak ada proteksi, input langsung diproses)",
        "  - Database: MySQL, Backend PHP 7.0"
    ])

    # Slide 5: Reconnaissance & Scanning
    add_bullet_slide(prs, "Reconnaissance & Scanning", [
        "Nmap (Port & Service Scan):",
        "  - 80/tcp open: Apache httpd 2.4.25",
        "  - 3306/tcp open: MySQL",
        "WhatWeb (Fingerprinting):",
        "  - Mengidentifikasi Apache, PHP, Cookies, Redirect.",
        "Nikto (Web Vulnerability Scanner):",
        "  - Apache/2.4.25 (versi lama).",
        "  - Cookie PHPSESSID tanpa flag HttpOnly (rentan XSS).",
        "Gobuster & OWASP ZAP:",
        "  - Directory Bruteforce & Automated Scanning."
    ])

    # Slide 6: Temuan 1 - SQL Injection
    add_bullet_slide(prs, "Temuan 1: SQL Injection", [
        "Deskripsi: Input pengguna dimasukkan ke query SQL tanpa sanitasi.",
        "Parameter Rentan: id (GET parameter).",
        "Eksploitasi: Payload `1' OR '1'='1` mengembalikan semua data dari database (Authentication Bypass / Data Extraction).",
        "Severity: CRITICAL (CVSS v4.0: 9.3).",
        "Mitigasi: Gunakan Prepared Statements / Parameterized Queries."
    ])

    # Slide 7: Temuan 2 - Reflected XSS
    add_bullet_slide(prs, "Temuan 2: Reflected XSS", [
        "Deskripsi: Payload JavaScript disisipkan via parameter URL dan direfleksikan ke halaman tanpa encoding.",
        "Parameter Rentan: name (GET parameter).",
        "Eksploitasi: Script berbahaya dieksekusi di browser korban (dapat mencuri cookie/session).",
        "Severity: MEDIUM (CVSS v4.0: 6.1).",
        "Mitigasi: Implementasikan Output Encoding dan validasi input di sisi server."
    ])

    # Slide 8: Temuan 3 - Stored XSS
    add_bullet_slide(prs, "Temuan 3: Stored XSS", [
        "Deskripsi: Payload disimpan permanen di database server dan dieksekusi otomatis setiap kali halaman dibuka.",
        "Fitur Rentan: Message textarea pada form Guestbook.",
        "Dampak: Mass Session Hijacking, Admin Account Takeover.",
        "Severity: HIGH (CVSS v4.0: 7.2).",
        "Mitigasi: Sanitasi input sebelum simpan, encode output saat ditampilkan, dan gunakan CSP (Content Security Policy)."
    ])

    # Slide 9: Temuan 4 - Unrestricted File Upload (RCE)
    add_bullet_slide(prs, "Temuan 4: Unrestricted File Upload (RCE)", [
        "Deskripsi: Aplikasi tidak memvalidasi jenis file yang diupload.",
        "Eksploitasi: Mengupload file PHP web shell (shell.php).",
        "Dampak: Remote Code Execution (RCE), penyerang dapat menjalankan perintah OS dan mengontrol server penuh.",
        "Severity: CRITICAL (CVSS v4.0: 9.3).",
        "Mitigasi: Whitelist ekstensi file, validasi MIME type, dan nonaktifkan eksekusi PHP di direktori uploads."
    ])

    # Slide 10: Ringkasan Penilaian Risiko
    add_bullet_slide(prs, "Ringkasan Penilaian Risiko", [
        "CRITICAL (CVSS 9.0+):",
        "  - SQL Injection (9.3)",
        "  - File Upload RCE (9.3)",
        "HIGH (CVSS 7.0 - 8.9):",
        "  - Stored XSS (7.2)",
        "MEDIUM (CVSS 4.0 - 6.9):",
        "  - Reflected XSS (6.1)",
        "  - Missing HttpOnly Cookie (5.4)",
        "  - Apache Outdated (5.0)",
        "  - MySQL Port Exposed (4.3)"
    ])

    # Slide 11: Rekomendasi Prioritas & Kesimpulan
    add_bullet_slide(prs, "Rekomendasi Prioritas & Kesimpulan", [
        "Prioritas 1 (Critical): Perbaiki SQL Injection (Prepared Statements) dan validasi ketat pada file upload.",
        "Prioritas 2 (High): Implementasikan output encoding & CSP untuk Stored XSS.",
        "Prioritas 3 (Medium): Update Apache dan batasi akses MySQL.",
        "Lessons Learned:",
        "  - Sanitasi input dan output adalah kewajiban developer.",
        "  - Vulnerability Assessment rutin diperlukan untuk menjaga keamanan sistem.",
        "  - Keamanan berlapis (Defense in Depth) sangat penting."
    ])

    # Save presentation
    prs.save("Presentasi_VA_DVWA.pptx")
    print("Presentasi berhasil dibuat: Presentasi_VA_DVWA.pptx")

if __name__ == "__main__":
    create_presentation()
